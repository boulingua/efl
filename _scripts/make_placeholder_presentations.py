"""Generate placeholder presentation .pptx files for every EFL Unit.

Reads `_resources/curriculum_outline.yml` and emits one one-slide .pptx per
Unit at:

    static/materials/presentations/<track>/kl<NN>/unit<NN>_<slug>.pptx

Mirrors `_scripts/make_placeholder_worksheets.py` in shape so the two
material tracks stay parallel. When real decks exist, drop them at the
same canonical path; no site-code change needed.

Each .pptx has a single 16:9 slide carrying the unit title, track,
Klassenstufe, Niveau, and a "Placeholder — replace with final deck"
caption. Authored by S. Le Boulanger; metadata is written into the
.pptx core properties so attribution travels with the file.
"""
from __future__ import annotations

import argparse
import pathlib
import sys

import yaml
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

HERE = pathlib.Path(__file__).resolve().parent
REPO = HERE.parent
OUTLINE = REPO / "_resources" / "curriculum_outline.yml"
DEFAULT_OUT = REPO / "static" / "materials" / "presentations"

AUTHOR = "S. Le Boulanger"


def render_one(out_path: pathlib.Path, track: str, klasse: int,
               niveau: str, unit_nr: int, slug: str, title: str) -> None:
    out_path.parent.mkdir(parents=True, exist_ok=True)

    pres = Presentation()
    pres.slide_width = Inches(13.333)
    pres.slide_height = Inches(7.5)

    blank_layout = pres.slide_layouts[6]
    slide = pres.slides.add_slide(blank_layout)

    # Title.
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.7),
                                  Inches(12), Inches(1.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"Unit {unit_nr}: {title}"
    p.runs[0].font.size = Pt(40)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = RGBColor(0x22, 0x22, 0x22)

    # Subtitle line.
    sb = slide.shapes.add_textbox(Inches(0.7), Inches(2.0),
                                  Inches(12), Inches(0.6))
    sp = sb.text_frame.paragraphs[0]
    track_label = "G+M" if track == "gm" else "E"
    sp.text = f"Track {track_label} · Klasse {klasse} · Niveau {niveau}"
    sp.runs[0].font.size = Pt(20)
    sp.runs[0].font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    # Placeholder note.
    nb = slide.shapes.add_textbox(Inches(0.7), Inches(3.6),
                                  Inches(12), Inches(2.4))
    np = nb.text_frame.paragraphs[0]
    np.text = "Placeholder — replace with final presentation."
    np.runs[0].font.size = Pt(24)
    np.runs[0].font.italic = True
    np.runs[0].font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    # Attribution footer.
    ab = slide.shapes.add_textbox(Inches(0.7), Inches(6.7),
                                  Inches(12), Inches(0.5))
    ap = ab.text_frame.paragraphs[0]
    ap.text = f"© {AUTHOR} · CC-BY-SA 4.0"
    ap.runs[0].font.size = Pt(12)
    ap.runs[0].font.color.rgb = RGBColor(0x88, 0x88, 0x88)

    # Core properties (metadata).
    cp = pres.core_properties
    cp.author = AUTHOR
    cp.title = f"Unit {unit_nr}: {title}"
    cp.subject = f"EFL · Track {track_label} · Klasse {klasse}"

    pres.save(str(out_path))


def iter_units(outline: dict):
    for course in outline.get("courses", []):
        track = course["track"]
        klasse = int(course["klassenstufe"])
        niveau = course.get("niveau", "")
        for unit in course.get("units", []):
            yield {
                "track": track,
                "klasse": klasse,
                "niveau": unit.get("niveau", niveau),
                "unit_nr": int(unit["unit_nr"]),
                "slug": unit["slug"],
                "title": unit["title"],
            }


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--out", default=str(DEFAULT_OUT))
    args = ap.parse_args()

    if not OUTLINE.exists():
        print(f"WARN: {OUTLINE} not found.", file=sys.stderr)
        return 0

    with OUTLINE.open(encoding="utf-8") as f:
        outline = yaml.safe_load(f) or {}

    out_base = pathlib.Path(args.out)
    count = 0
    for u in iter_units(outline):
        nn = f"{u['unit_nr']:02d}"
        kk = f"{u['klasse']:02d}"
        path = out_base / u["track"] / f"kl{kk}" / f"unit{nn}_{u['slug']}.pptx"
        render_one(path, u["track"], u["klasse"], u["niveau"],
                   u["unit_nr"], u["slug"], u["title"])
        count += 1

    print(f"Wrote {count} placeholder .pptx file(s) under {out_base}/")
    return 0


if __name__ == "__main__":
    sys.exit(main())
