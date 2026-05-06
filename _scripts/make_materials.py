"""Generate per-unit Materials placeholders + inject frontmatter.

Modelled on `boulingua/fle`'s consolidated generator: one script
produces the .pptx, the .pdf, both thumbnail PNGs, and rewrites the
unit page's `presentation:` + `worksheet:` frontmatter blocks. Idempotent.

EFL adaptation:
  - Source pages live in page bundles
    (`content/track-{e,gm}/kl<NN>/units/unit<NN>-<slug>/index.md`),
    so the FLE-style "filename stem" doesn't carry course provenance.
    We derive a flat slug ourselves from the path:
      track-{e,gm}_kl<NN>_unit<NN>-<slug>
  - Output binaries land flat under
    `static/materials/presentations/<flat-slug>.{pptx,png}` and
    `static/materials/worksheets/<flat-slug>.{pdf,png}` — same structure
    as FLE so manual replacement is a one-file drop.
  - URLs in frontmatter are hardcoded with the `/efl/` prefix because
    Hugo's `relURL` does not prepend baseURL to leading-slash inputs.
    If the site ever moves off `/efl/`, regenerate.

Real thumbnails:
  - PDFs render via pypdfium2 (no system deps).
  - PPTXs render via LibreOffice headless → PDF → pypdfium2.
  - Both fall back to a PIL synthetic title card when the toolchain is
    missing (typical local dev on Windows). CI always uses real renders.
"""
from __future__ import annotations

import io
import pathlib
import re
import shutil
import subprocess
import sys
import tempfile
from typing import Optional

REPO = pathlib.Path(__file__).resolve().parent.parent
CONTENT = REPO / "content"
STATIC_BASE = REPO / "static" / "materials"
PRES_DIR = STATIC_BASE / "presentations"
WORK_DIR = STATIC_BASE / "worksheets"

URL_PREFIX = "/efl"  # bake the GitHub-Pages basePath; Hugo's relURL won't.
AUTHOR = "S. Le Boulanger"

YAML_RE = re.compile(r"^---\n(.*?)\n---\n", re.DOTALL)


# ----- Path / slug helpers -----------------------------------------------

UNIT_BUNDLE_RE = re.compile(
    r"^content/track-(e|gm)/kl(\d{2})/units/(unit\d{2}-[^/]+)/index\.md$"
)


def list_unit_pages() -> list[pathlib.Path]:
    out = []
    for p in CONTENT.rglob("units/*/index.md"):
        rel = p.relative_to(REPO).as_posix()
        if not UNIT_BUNDLE_RE.match(rel):
            continue
        if p.parent.name.endswith("-exam"):
            continue
        out.append(p)
    return sorted(out)


def flat_slug(md_path: pathlib.Path) -> str:
    """Derive a flat slug from a unit-bundle path:
       content/track-e/kl06/units/unit01-foo/index.md -> track-e_kl06_unit01-foo"""
    rel = md_path.relative_to(REPO).as_posix()
    m = UNIT_BUNDLE_RE.match(rel)
    if not m:
        raise ValueError(f"unexpected unit path: {rel}")
    track, kl, unit_dir = m.groups()
    return f"track-{track}_kl{kl}_{unit_dir}"


# ----- Frontmatter -------------------------------------------------------

def parse_fm(text: str) -> tuple[dict, str, str]:
    m = YAML_RE.match(text)
    if not m:
        return {}, "", text
    fm_text = m.group(1)
    fm: dict = {}
    for line in fm_text.splitlines():
        sm = re.match(r"^([A-Za-z][\w-]*):\s*(.*)$", line)
        if sm:
            key, val = sm.group(1), sm.group(2).strip()
            if val and not val.startswith("|") and not val.startswith(">"):
                fm[key] = val.strip().strip('"')
    return fm, fm_text, text[m.end():]


def upsert_fm_block(fm_text: str, key: str, lines: list[str]) -> str:
    pattern = re.compile(
        rf"^{re.escape(key)}:\s*\n(?:[ \t]+\S.*\n?)*",
        re.MULTILINE,
    )
    block = key + ":\n" + "\n".join("  " + l for l in lines) + "\n"
    if pattern.search(fm_text):
        return pattern.sub(block, fm_text, count=1)
    if not fm_text.endswith("\n"):
        fm_text += "\n"
    return fm_text + block


# ----- Binary generation -------------------------------------------------

def make_pptx(out: pathlib.Path, title: str, subtitle: str) -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)

    tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.7), Inches(12), Inches(1.4))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.runs[0].font.size = Pt(40)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = RGBColor(0x22, 0x22, 0x22)

    sb = slide.shapes.add_textbox(Inches(0.7), Inches(2.2), Inches(12), Inches(0.6))
    sp = sb.text_frame.paragraphs[0]
    sp.text = subtitle
    sp.runs[0].font.size = Pt(20)
    sp.runs[0].font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    nb = slide.shapes.add_textbox(Inches(0.7), Inches(3.6), Inches(12), Inches(2.4))
    np = nb.text_frame.paragraphs[0]
    np.text = "Placeholder — replace with final presentation."
    np.runs[0].font.size = Pt(24)
    np.runs[0].font.italic = True
    np.runs[0].font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    ab = slide.shapes.add_textbox(Inches(0.7), Inches(6.7), Inches(12), Inches(0.5))
    ap = ab.text_frame.paragraphs[0]
    ap.text = f"© {AUTHOR} · CC-BY-SA 4.0"
    ap.runs[0].font.size = Pt(12)
    ap.runs[0].font.color.rgb = RGBColor(0x88, 0x88, 0x88)

    cp = prs.core_properties
    cp.author = AUTHOR
    cp.title = title

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))


def make_pdf(out: pathlib.Path, title: str, subtitle: str) -> None:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib.units import cm
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer

    out.parent.mkdir(parents=True, exist_ok=True)
    doc = SimpleDocTemplate(
        str(out), pagesize=A4,
        leftMargin=2.5 * cm, rightMargin=2.5 * cm,
        topMargin=2.5 * cm, bottomMargin=2.5 * cm,
        author=AUTHOR, title=title,
    )
    styles = getSampleStyleSheet()
    story = [
        Paragraph(title, styles["Title"]),
        Spacer(1, 0.5 * cm),
        Paragraph(subtitle, styles["Heading3"]),
        Spacer(1, 1 * cm),
        Paragraph(
            "Worksheet — placeholder. Will be replaced with the final version.",
            styles["BodyText"],
        ),
        Spacer(1, 4 * cm),
        Paragraph(
            f"<i>© {AUTHOR} · CC-BY-SA 4.0</i>",
            styles["BodyText"],
        ),
    ]
    doc.build(story)


# ----- Thumbnails --------------------------------------------------------

def find_libreoffice() -> str | None:
    for cmd in ("soffice", "libreoffice"):
        if shutil.which(cmd):
            return cmd
    for p in (
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
    ):
        if pathlib.Path(p).is_file():
            return p
    return None


SOFFICE = find_libreoffice()


def render_pdf_first_page(pdf: pathlib.Path, png: pathlib.Path,
                          width: int = 800) -> bool:
    try:
        import pypdfium2 as pdfium
    except ImportError:
        return False
    try:
        doc = pdfium.PdfDocument(str(pdf))
        if len(doc) == 0:
            return False
        page = doc[0]
        scale = width / page.get_width()
        img = page.render(scale=scale).to_pil()
        png.parent.mkdir(parents=True, exist_ok=True)
        img.save(png, format="PNG", optimize=True)
        return True
    except Exception as e:
        print(f"  pypdfium2 fail on {pdf.name}: {e}", file=sys.stderr)
        return False


def render_pptx_first_slide(pptx: pathlib.Path, png: pathlib.Path) -> bool:
    if SOFFICE is None:
        return False
    with tempfile.TemporaryDirectory() as td:
        td_path = pathlib.Path(td)
        try:
            r = subprocess.run(
                [SOFFICE, "--headless", "--convert-to", "pdf",
                 "--outdir", str(td_path), str(pptx)],
                capture_output=True, text=True, timeout=120,
            )
        except subprocess.TimeoutExpired:
            return False
        if r.returncode != 0:
            return False
        produced = td_path / (pptx.stem + ".pdf")
        if not produced.is_file():
            return False
        return render_pdf_first_page(produced, png)


def make_thumbnail(png: pathlib.Path, title: str, kind: str,
                   accent_hex: str) -> None:
    """PIL synthetic title-card fallback when real render is unavailable."""
    from PIL import Image, ImageDraw, ImageFont
    w, h = 800, 450
    bg = (245, 247, 250)
    fg = (34, 34, 34)
    accent = tuple(int(accent_hex[i:i+2], 16) for i in (1, 3, 5))
    img = Image.new("RGB", (w, h), bg)
    draw = ImageDraw.Draw(img)
    draw.rectangle([(0, 0), (w, 8)], fill=accent)
    try:
        font = ImageFont.truetype("DejaVuSans.ttf", 30)
        small = ImageFont.truetype("DejaVuSans.ttf", 16)
    except (OSError, IOError):
        font = ImageFont.load_default()
        small = font
    draw.text((40, 60), kind.upper(), font=small, fill=accent)
    draw.text((40, 96), title[:60], font=font, fill=fg)
    draw.text((40, h - 60),
              "Synthetic preview — install LibreOffice for real PPTX renders.",
              font=small, fill=(140, 140, 140))
    png.parent.mkdir(parents=True, exist_ok=True)
    img.save(png, format="PNG", optimize=True)


# ----- Per-unit driver ---------------------------------------------------

def short_label(fm: dict) -> str:
    track = fm.get("track", "")
    klasse = fm.get("klassenstufe", "")
    niveau = fm.get("niveau", "")
    parts = []
    if track:
        parts.append(f"Track {'G+M' if track == 'gm' else 'E'}")
    if klasse:
        parts.append(f"Klasse {klasse}")
    if niveau:
        parts.append(f"Niveau {niveau}")
    return " · ".join(parts) or "EFL"


def process_unit(md_path: pathlib.Path, force: bool = False) -> dict:
    text = md_path.read_text(encoding="utf-8")
    fm, fm_text, body = parse_fm(text)
    slug = flat_slug(md_path)
    title = fm.get("title", slug)
    subtitle = short_label(fm)

    pres_pptx = PRES_DIR / f"{slug}.pptx"
    pres_png = PRES_DIR / f"{slug}.png"
    work_pdf = WORK_DIR / f"{slug}.pdf"
    work_png = WORK_DIR / f"{slug}.png"

    if force or not pres_pptx.is_file():
        make_pptx(pres_pptx, title, subtitle)
    if force or not work_pdf.is_file():
        make_pdf(work_pdf, title, subtitle)

    # Always re-render thumbnails — they reflect whatever binary is at
    # the path right now, including manually-dropped real materials.
    if not render_pptx_first_slide(pres_pptx, pres_png):
        make_thumbnail(pres_png, title, "Presentation", "#1A73E8")
    if not render_pdf_first_page(work_pdf, work_png):
        make_thumbnail(work_png, title, "Worksheet", "#2E7D32")

    pres_url = f"{URL_PREFIX}/materials/presentations/{slug}.pptx"
    pres_thumb_url = f"{URL_PREFIX}/materials/presentations/{slug}.png"
    work_url = f"{URL_PREFIX}/materials/worksheets/{slug}.pdf"
    work_thumb_url = f"{URL_PREFIX}/materials/worksheets/{slug}.png"

    new_fm = upsert_fm_block(fm_text, "presentation", [
        f'file: "{pres_url}"',
        f'thumbnail: "{pres_thumb_url}"',
    ])
    new_fm = upsert_fm_block(new_fm, "worksheet", [
        f'file: "{work_url}"',
        f'thumbnail: "{work_thumb_url}"',
    ])

    new_text = "---\n" + new_fm.rstrip("\n") + "\n---\n" + body
    if new_text != text:
        md_path.write_text(new_text, encoding="utf-8")
    return {"slug": slug, "title": title}


def main() -> int:
    import argparse
    ap = argparse.ArgumentParser()
    ap.add_argument("--force", action="store_true",
                    help="Overwrite existing .pptx / .pdf even when present "
                         "(default skips them so manually-placed real "
                         "materials survive). Thumbnails always re-render.")
    args = ap.parse_args()

    units = list_unit_pages()
    if not units:
        print("no unit pages found", file=sys.stderr)
        return 1
    PRES_DIR.mkdir(parents=True, exist_ok=True)
    WORK_DIR.mkdir(parents=True, exist_ok=True)

    print(f"LibreOffice: {SOFFICE or 'NOT FOUND (synthetic .pptx thumbs)'}")
    for i, p in enumerate(units, 1):
        info = process_unit(p, force=args.force)
        if i % 30 == 0 or i == len(units):
            print(f"  [{i}/{len(units)}] {info['slug']}")
    print(f"Done — {len(units)} units; "
          f"{len(list(PRES_DIR.glob('*.pptx')))} pptx, "
          f"{len(list(WORK_DIR.glob('*.pdf')))} pdf, "
          f"{len(list(PRES_DIR.glob('*.png')))} pres-thumbs, "
          f"{len(list(WORK_DIR.glob('*.png')))} work-thumbs.")
    return 0


if __name__ == "__main__":
    sys.exit(main())
